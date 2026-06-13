"""Monta o .pptx de uma aula a partir do template `aulas/templates/MX AY.pptx`.

Regras (definidas pelo coordenador):
- Slide 1 (capa): troca apenas número/nome do módulo e número/nome da aula.
- Slide 2 (modelo): caixa de texto à direita. Cada bloco do texto (separado
  por linhas de hífens) vira uma cópia desse slide com o bloco colado.
- Não altera design/layout do template; não altera o texto da aula.
"""

from __future__ import annotations

import copy
import io
import re
from pathlib import Path

from pptx import Presentation
from pptx.oxml.ns import qn

# Linha separadora de blocos: linha contendo apenas 3+ hifens. O texto colado do
# NotebookLM separa os parágrafos com "---" (3 hifens), então exigir 4+ fazia
# nenhuma quebra casar e tudo caía num único slide.
_BLOCK_SEP = re.compile(r"(?m)^[ \t]*-{3,}[ \t]*$")

# Separador de parágrafos (linha em branco). Usado como fallback quando o texto
# do NotebookLM vem SEM os "---": aí cada parágrafo (bloco separado por linha em
# branco) vira um slide.
_PARA_SEP = re.compile(r"\n[ \t]*\n+")

# Citações do NotebookLM no corpo do texto: números entre colchetes, incluindo
# listas/intervalos como [1], [1, 2], [3-5]. Removidas do PPTX (a rastreabilidade
# fica nos .md de bibliografia e no slide final de referências). NÃO casa links
# markdown [texto](url) nem [sic], porque exige só dígitos/vírgulas/traços dentro.
_CITATION_RE = re.compile(r"[ \t]*\[\d+(?:\s*[,–—-]\s*\d+)*\]")

TEMPLATE_FILENAME = "MX AY.pptx"

# Tamanho da fonte (pt) do texto dos slides de conteúdo. O template traz
# 18pt, grande demais para blocos de 600-800 caracteres — 14pt evita que a
# caixa de texto ultrapasse os limites do slide.
CONTENT_FONT_SIZE_PT = 14

# Tokens do template da capa -> valor a inserir (preenchido em runtime).
_CAPA_TOKENS = {
    "modulo_num": "Módulo X",
    "modulo_nome": "[INSERIR NOME DO MÓDULO AQUI]",
    "aula_num": "[insira número da aula aqui]",
    "aula_nome": "[insira nome da aula aqui]",
}


def template_path() -> Path:
    """Resolve o caminho do template, tanto no repo quanto no container."""
    from prompt_loader import TEMPLATE_DIR_CANDIDATES

    for base in TEMPLATE_DIR_CANDIDATES:
        candidate = base / TEMPLATE_FILENAME
        if candidate.is_file():
            return candidate
    raise FileNotFoundError(
        f"Template '{TEMPLATE_FILENAME}' não encontrado em aulas/templates."
    )


def strip_citations(texto: str) -> str:
    """Remove citações no formato [N] (ex.: [1], [1, 2], [3-5]) do corpo do texto.
    Tira também o espaço que sobraria antes da pontuação e colapsa espaços duplos.
    Preserva links markdown [texto](url) e marcadores como [sic]."""
    out = _CITATION_RE.sub("", texto or "")
    out = re.sub(r" +([.,;:!?])", r"\1", out)  # "palavra ." -> "palavra."
    out = re.sub(r"[ \t]{2,}", " ", out)
    return out


def split_blocks(texto: str) -> list[str]:
    """Divide o texto da aula em blocos (1 por slide).

    Prioriza as linhas separadoras de hífens ("---"). Se o texto NÃO tiver
    nenhuma "---" (o NotebookLM às vezes cola sem elas, só com parágrafos),
    cai para dividir por parágrafo: cada bloco separado por linha em branco
    vira um slide."""
    texto = texto or ""
    if _BLOCK_SEP.search(texto):
        parts = _BLOCK_SEP.split(texto)
    else:
        parts = _PARA_SEP.split(texto)
    return [b.strip() for b in parts if b.strip()]


def _capa_title_shape(slide):
    """Localiza o placeholder de título da capa (idx=0)."""
    for shape in slide.shapes:
        try:
            if shape.is_placeholder and shape.placeholder_format.idx == 0:
                return shape
        except Exception:
            pass
    for shape in slide.shapes:
        if shape.has_text_frame and _CAPA_TOKENS["modulo_nome"] in shape.text_frame.text:
            return shape
    return slide.shapes[0]


def _fill_capa(slide, modulo_num, modulo_nome: str, aula_num, aula_nome: str) -> None:
    """Troca os tokens da capa preservando a formatação de cada run."""
    repl = {
        _CAPA_TOKENS["modulo_num"]: f"Módulo {modulo_num}",
        _CAPA_TOKENS["modulo_nome"]: str(modulo_nome).strip(),
        _CAPA_TOKENS["aula_num"]: str(aula_num).strip(),
        _CAPA_TOKENS["aula_nome"]: str(aula_nome).strip(),
    }
    title = _capa_title_shape(slide)
    for para in title.text_frame.paragraphs:
        for run in para.runs:
            novo = run.text
            for old, new in repl.items():
                if old in novo:
                    novo = novo.replace(old, new)
            if novo != run.text:
                run.text = novo


def _content_textbox(slide):
    """Localiza a caixa de texto do slide de conteúdo."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            return shape
    raise ValueError("Slide de conteúdo do template não tem caixa de texto.")


def _fill_textbox(
    slide,
    block: str,
    font_size_pt: int = CONTENT_FONT_SIZE_PT,
    align: str | None = None,
) -> None:
    """Cola um bloco de texto na caixa, preservando o parágrafo/run modelo
    do template e ajustando o tamanho da fonte. Linhas em branco viram
    parágrafos vazios com a mesma altura de fonte (espaçamento).
    `align`: se informado (ex.: "l"), sobrescreve o alinhamento do template."""
    shape = _content_textbox(slide)
    txBody = shape.text_frame._txBody
    paragraphs = txBody.findall(qn("a:p"))
    proto = copy.deepcopy(paragraphs[0])
    for para in paragraphs:
        txBody.remove(para)

    # OOXML expressa o tamanho da fonte em centésimos de ponto.
    sz = str(int(round(font_size_pt * 100)))
    for linha in block.split("\n"):
        novo = copy.deepcopy(proto)
        runs = novo.findall(qn("a:r"))
        if align is not None:
            p_pr = novo.find(qn("a:pPr"))
            if p_pr is not None:
                p_pr.set("algn", align)
        # endParaRPr define a altura de parágrafos vazios; alinha à fonte
        # (sem isso, linhas em branco herdariam os 28pt do template).
        end_pr = novo.find(qn("a:endParaRPr"))
        if end_pr is not None:
            end_pr.set("sz", sz)
        if linha.strip() == "":
            # Parágrafo vazio: remove os runs, mantém pPr/endParaRPr.
            for run in runs:
                novo.remove(run)
        else:
            # Mantém só o primeiro run (formatação modelo) e ajusta texto/tamanho.
            for run in runs[1:]:
                novo.remove(run)
            runs[0].find(qn("a:t")).text = linha
            r_pr = runs[0].find(qn("a:rPr"))
            if r_pr is not None:
                r_pr.set("sz", sz)
        txBody.append(novo)


def _duplicate_slide(prs, src_slide):
    """Cria uma nova cópia de `src_slide` no fim da apresentação."""
    dst = prs.slides.add_slide(src_slide.slide_layout)
    # Remove os placeholders herdados do layout.
    for shape in list(dst.shapes):
        shape._element.getparent().remove(shape._element)
    # Copia os shapes do slide de origem.
    for shape in src_slide.shapes:
        dst.shapes._spTree.append(copy.deepcopy(shape._element))
    return dst


# ----------------------------------------------------------------------
# Slide final de referências (compilado dos .md de bibliografia curados)
# ----------------------------------------------------------------------

_LINK_RE = re.compile(r"\[([^\]\n]+)\]\((https?://[^)\s]+)\)")
_BOLD_RE = re.compile(r"\*\*([^*]+)\*\*")

# Tamanhos candidatos para a fonte do slide de referências (maior -> menor).
_REF_FONT_SIZES = (12, 11, 10, 9)

# Rótulo amigável dos livros-base.
_LIVRO_NOMES = {
    "tratado": "Tratado de Ginecologia (FEBRASGO)",
    "williams": "Williams Ginecologia",
}


def _link_entries(md: str) -> list[dict]:
    """Extrai de um .md de bibliografia as linhas que contêm link markdown.
    Linhas sem link (cabeçalhos, lacunas, observações) são ignoradas."""
    entries: list[dict] = []
    for line in (md or "").splitlines():
        m = _LINK_RE.search(line)
        if not m:
            continue
        bold = _BOLD_RE.search(line[: m.start()])
        meta = line[m.end():].strip().lstrip("—–-· ").strip()
        entries.append(
            {
                "fonte": bold.group(1).strip() if bold else "",
                "titulo": m.group(1).strip(),
                "url": m.group(2).strip(),
                "meta": meta,
            }
        )
    return entries


def _livro_entries(md: str) -> list[str]:
    """Extrai 'Livro — capítulo' de capitulos_livros.md (o link da extração
    fica no fim da linha e não entra no slide)."""
    out: list[str] = []
    for line in (md or "").splitlines():
        if not _LINK_RE.search(line):
            continue
        bold = _BOLD_RE.search(line)
        if not bold:
            continue
        livro = _LIVRO_NOMES.get(bold.group(1).strip().lower(), bold.group(1).strip())
        resto = line[bold.end():].lstrip(" —–-")
        for marca in ("· confiança", "·confiança", "→"):
            i = resto.find(marca)
            if i != -1:
                resto = resto[:i]
        cap = re.sub(r"^capítulo:\s*", "cap. ", resto.strip().rstrip("·").strip())
        out.append(f"{livro} — {cap}" if cap else livro)
    return out


def compor_referencias(arquivos: dict) -> str:
    """Compila o texto do slide final de referências a partir dos .md
    curados (`diretrizes_consensos.md`, `pubmed_busca.md`, `uptodate.md`,
    `capitulos_livros.md`). Retorna "" se não houver nada."""
    secoes: list[str] = []

    diretrizes = _link_entries(arquivos.get("diretrizes_consensos.md", ""))
    if diretrizes:
        linhas = ["Diretrizes e consensos"]
        for e in diretrizes:
            cab = f"• {e['fonte']} — {e['titulo']}" if e["fonte"] else f"• {e['titulo']}"
            linhas += [cab, f"  {e['url']}"]
        secoes.append("\n".join(linhas))

    pubmed = _link_entries(arquivos.get("pubmed_busca.md", ""))
    if pubmed:
        linhas = ["Artigos (PubMed)"]
        for e in pubmed:
            extra = f" ({e['meta']})" if e["meta"] else ""
            linhas += [f"• {e['titulo']}{extra}", f"  {e['url']}"]
        secoes.append("\n".join(linhas))

    uptodate = _link_entries(arquivos.get("uptodate.md", ""))
    if uptodate:
        linhas = ["UpToDate"]
        for e in uptodate:
            linhas += [f"• {e['titulo']}", f"  {e['url']}"]
        secoes.append("\n".join(linhas))

    livros = _livro_entries(arquivos.get("capitulos_livros.md", ""))
    if livros:
        linhas = ["Livros e capítulos"]
        linhas += [f"• {l}" for l in livros]
        secoes.append("\n".join(linhas))

    if not secoes:
        return ""
    return "Referências\n\n" + "\n\n".join(secoes)


def _fit_ref_font(text: str, box_width_emu: int, usable_height_in: float = 6.4) -> int:
    """Escolhe o maior tamanho de fonte (de `_REF_FONT_SIZES`) que faz o
    slide de referências caber na altura util do slide."""
    width_pt = box_width_emu / 914400 * 72 - 14  # desconta as margens internas
    for size in _REF_FONT_SIZES:
        chars_per_line = max(20, int(width_pt / (0.5 * size)))
        linhas = 0
        for ln in text.split("\n"):
            linhas += max(1, -(-len(ln) // chars_per_line))  # ceil
        if linhas * 1.2 * size / 72 <= usable_height_in:
            return size
    return _REF_FONT_SIZES[-1]


def build_pptx(
    texto: str,
    modulo_num,
    modulo_nome: str,
    aula_num,
    aula_nome: str,
    template_file: Path | None = None,
    content_font_size_pt: int = CONTENT_FONT_SIZE_PT,
    referencias_text: str | None = None,
) -> tuple[bytes, int]:
    """Monta o .pptx da aula. Retorna (bytes, numero_de_slides).
    Se `referencias_text` for informado, adiciona um slide final com as
    referências (fonte reduzida automaticamente para caber)."""
    # Remove as citações [N] do corpo antes de dividir em slides (a rastreabilidade
    # fica nos .md de bibliografia e no slide final de referências).
    texto = strip_citations(texto)
    blocks = split_blocks(texto)
    if not blocks:
        blocks = [(texto or "").strip() or "(sem texto)"]

    referencias = (referencias_text or "").strip()
    # Cópias do slide-modelo necessárias: 1 por bloco + 1 de referências.
    n_slides_conteudo = len(blocks) + (1 if referencias else 0)

    prs = Presentation(str(template_file or template_path()))
    if len(prs.slides) < 2:
        raise ValueError("Template inválido: esperados 2 slides (capa + conteúdo).")

    _fill_capa(prs.slides[0], modulo_num, modulo_nome, aula_num, aula_nome)

    content_template = prs.slides[1]
    # Cria as cópias necessárias ANTES de preencher (mantém o modelo intacto).
    for _ in range(n_slides_conteudo - 1):
        _duplicate_slide(prs, content_template)

    # Slides de conteúdo ocupam os índices 1..N.
    for idx, block in enumerate(blocks):
        _fill_textbox(prs.slides[idx + 1], block, content_font_size_pt)

    # Slide final de referências (índice N+1), fonte ajustada para caber.
    if referencias:
        ref_slide = prs.slides[len(blocks) + 1]
        box_width = _content_textbox(ref_slide).width
        ref_size = _fit_ref_font(referencias, box_width)
        _fill_textbox(ref_slide, referencias, ref_size, align="l")

    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue(), len(prs.slides)
